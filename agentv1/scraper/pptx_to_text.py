# pptx_to_text.py
# pip install python-pptx

import io, os, re, datetime
from pptx import Presentation

# ---------- Extraction ----------

def _shape_text(shape) -> list[str]:
    """Extract visible text from a single shape (including grouped shapes + tables)."""
    lines = []

    # Grouped shapes
    if hasattr(shape, "shapes"):  # group shape
        for shp in shape.shapes:
            lines.extend(_shape_text(shp))
        return lines

    # Tables
    if getattr(shape, "has_table", False):
        tbl = shape.table
        for r in tbl.rows:
            row_cells = []
            for c in r.cells:
                row_cells.append(c.text.replace("\r", "\n").strip())
            # join cells with tabs (keeps structure but stays text-only)
            lines.append("\t".join(filter(None, row_cells)))
        return lines

    # Text frames / placeholders
    if hasattr(shape, "has_text_frame") and shape.has_text_frame:
        tf = shape.text_frame
        for p in tf.paragraphs:
            txt = "".join(run.text for run in p.runs)
            lvl = getattr(p, "level", 0) or 0
            txt = txt.replace("\r", "\n").strip()
            if not txt:
                continue
            # basic bullet-style prefix based on indent level
            prefix = ("- " if p.level == 0 else "  " * p.level + "- ")
            lines.append(prefix + txt)
    return lines


def _slide_text(slide) -> list[str]:
    lines = []
    # Title first (if any)
    if getattr(slide, "shapes", None):
        for shp in slide.shapes:
            if shp.name and "Title" in shp.name and getattr(shp, "has_text_frame", False):
                title = shp.text.strip()
                if title:
                    lines.append(f"# {title}")

        # Then all shapes
        for shp in slide.shapes:
            lines.extend(_shape_text(shp))

    # Notes (speaker notes)
    if getattr(slide, "has_notes_slide", False) and slide.notes_slide:
        note = slide.notes_slide.notes_text_frame
        if note:
            note_txt = "\n".join(p.text for p in note.paragraphs).strip()
            if note_txt:
                lines.append("\n[Notes]\n" + note_txt)

    # Filter duplicated empties and strip
    flat = []
    for ln in lines:
        ln = (ln or "").strip()
        if ln:
            flat.append(ln)
    return flat


def pptx_bytes_to_text(buf: io.BytesIO) -> str:
    """Extract text from PPTX (in-memory) with python-pptx."""
    buf.seek(0)
    prs = Presentation(buf)
    slides_out = []
    for idx, s in enumerate(prs.slides, start=1):
        body = _slide_text(s)
        if body:
            slides_out.append(f"\n--- Slide {idx} ---\n" + "\n".join(body))
        else:
            slides_out.append(f"\n--- Slide {idx} ---\n")
    return "\n".join(slides_out).strip()


# ---------- Cleaning (reuse your style) ----------

_LIGATURE_MAP = {
    "ﬀ": "ff", "ﬁ": "fi", "ﬂ": "fl", "ﬃ": "ffi", "ﬄ": "ffl",
    "–": "-", "—": "-", "−": "-", "“": '"', "”": '"', "’": "'", "‘": "'",
    "•": "-", "·": "·"
}

def clean_for_embeddings(raw: str) -> str:
    t = raw or ""

    # Fix common ligatures / punctuation first
    for k, v in _LIGATURE_MAP.items():
        t = t.replace(k, v)

    # Normalize line endings
    t = t.replace("\r\n", "\n").replace("\r", "\n")

    # Remove isolated slide numbers like lines that are just digits
    t = re.sub(r"\n\s*\d+\s*\n", "\n", t)

    # Join obvious hard wraps within paragraphs (keep blank lines as paragraph breaks)
    t = re.sub(r"(?<![.!?;:])\n(?!\n|- )", " ", t)  # don’t join lines that look like list items

    # Collapse 3+ newlines → 2
    t = re.sub(r"\n{3,}", "\n\n", t)

    # Trim repeated spaces
    t = re.sub(r"[ \t]{2,}", " ", t)

    return t.strip()


# ---------- Save helper (mirrors your PDF version) ----------

def save_pptx_bytes_as_txt(buf: io.BytesIO, pptx_filename: str, output_dir: str = ".",
                           include_header: bool = True) -> str:
    """
    Convert PPTX bytes to cleaned .txt for embeddings.
    """
    base = os.path.splitext(os.path.basename(pptx_filename))[0]
    txt_path = os.path.join(output_dir, f"{base}.txt")

    text = pptx_bytes_to_text(buf)
    text = clean_for_embeddings(text)

    if include_header:
        today = datetime.date.today().isoformat()
        header = f"{base}\nDate: {today}\n\n"
        text = header + text

    with open(txt_path, "w", encoding="utf-8") as f:
        f.write(text)
    return txt_path


# ---------- Example usage ----------

# if __name__ == "__main__":
#     # Load a .pptx into memory (no temp files)
#     with open("L07 SQL Part 2.pptx", "rb") as f:
#         pptx_buf = io.BytesIO(f.read())

#     # Preview
#     preview = pptx_bytes_to_text(pptx_buf)
#     print(preview[:2000])

#     # Save to .txt (named after the .pptx)
#     out_path = save_pptx_bytes_as_txt(pptx_buf, "L07 SQL Part 2.pptx")
#     print(f"Saved text file at: {out_path}")
